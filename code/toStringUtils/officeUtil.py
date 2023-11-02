import queue
from docx import Document
from docx.oxml import OxmlElement
import aspose.words as aw
import re
import xlrd
from pptx import Presentation
import aspose.slides as slides
import os
import docx
import fitz
from datetime import datetime
import textract
from informationEngine.table_extract import XlsxDevider
from util import globalVar
from toStringUtils.picUtil import *
from util.extractInfo import *
from toStringUtils.universalUtil import *
import pandas as pd


"""
officeUtil: 解析 docx/pdf/wps/et
"""


# 日志模块
from util.logUtils import LoggerSingleton
TAG = "toStringUtils.officeUtil.py-"
logger = LoggerSingleton().get_logger()


######################### pdf file ########################################


# 提取pdf文档中的文本和图片
def pdf_file(file_path):

    pdf_file_name = file_path.split("/")[-1]
    result_image_path = "../workspace/image/pdf/"+pdf_file_name+'/'
    os.makedirs(result_image_path, exist_ok=True)

    try:
        pdf = fitz.open(file_path)
        text = []

        count = 1

        for itm, page in enumerate(pdf):
            try:
                tupleImage = page.get_images()
                text.append(page.get_text("text"))
                for xref0 in tupleImage:
                    xref = xref0[0]
                    img = pdf.extract_image(xref)
                    ext = img['ext']

                    imageFilename = os.path.join(
                        result_image_path, f"{os.path.basename(file_path)}_image{count}.png")
                    imgout = open(imageFilename, 'wb')
                    imgout.write(img["image"])
                    imgout.close()
                    text.append("$PictureIsHere$")
                    count += 1
            except:
                continue
        pdf.close()

        logger.info(TAG+"pdf_file()-文本信息:")
        logger.info(text)

        # 解析图片信息
        image_all_text_res = []
        if globalVar.flag_list[0] == True:
            image_all_text = ocr_table_batch(result_image_path)
            # image_all_text = ocr_batch_textract(image_dir)
            logger.info(TAG+"pdf_file()-image_all_text: ")
            for row in image_all_text:
                # logger.info(row)
                try:
                    single_result = " ".join(
                        [element for sublist in row[1:] for element in sublist])
                    image_all_text_res.append(single_result)
                except IndexError as e:
                    # 处理 IndexError 异常
                    logger.error(e)
            logger.info(TAG+"pdf_file()-图片文本信息:")
            logger.info(image_all_text_res)

            res = replace_picture_texts(text, image_all_text_res)
            return "\n".join(res)
        else:
            logger.info(TAG+"pdf_file()-不处理文件内部的图片!!")
            text = "\n".join(text).replace("$PictureIsHere$", "")
            return text

    except Exception as e:
        logger.error(e)
        return ""


######################### office/wps file ########################################


# 1-1.解析.doc/.wps/.docx
def docs_file(file_path, type):
    time = datetime.now().strftime("%Y%m%d%H%M%S%f")
    if type == ".doc":
        logger.info(TAG+"docs_file(): .doc")
        doc_file_dir = "../workspace/office/doc/"
        result_image_path = "../workspace/image/office/doc"
        docx_path = doc_file_dir + \
            time+"_"+file_path.replace(".doc", ".docx").split("/")[-1]
    elif type == ".wps":
        logger.info(TAG+"docs_file(): .wps")
        doc_file_dir = "../workspace/wps/wps/"
        result_image_path = "../workspace/image/wps/wps"
        docx_path = doc_file_dir + \
            time+"_"+file_path.replace(".doc", ".docx").split("/")[-1]
    elif type == ".docx":
        logger.info(TAG+"docs_file(): .docx")
        result_image_path = "../workspace/image/office/docx"
        image_dir = f"{result_image_path}/{time}_{os.path.basename(file_path)}/"
        return docx_file_info_extract(file_path, image_dir)
    else:
        return ""

    os.makedirs(doc_file_dir, exist_ok=True)
    logger.info(TAG+"docs_file(): "+docx_path)

    os.makedirs(result_image_path, exist_ok=True)
    doc_docx_name = time+"_"+file_path.split("/")[-1]

    # 使用Aspose.Words将.doc转换为.docx
    try:
        doc = aw.Document(file_path)
        doc.save(docx_path, aw.SaveFormat.DOCX)
        logger.info(TAG+"docs_file(): "+"Document conversion successful!")

        image_dir = f"{result_image_path}/{doc_docx_name}/"

        return docx_file_info_extract(docx_path, image_dir)
    except Exception as e:
        logger.error(e)
        return ""


# 1-2.提取.docx中的文本和图片文本信息
def docx_file_info_extract(docx_path, image_dir):
    logger.info(TAG+"docx_file_info_extract():")

    docx_text = " "

    target_docx = docx.Document(docx_path)

    # 是否成功设置图片占位符
    flag = False

    try:
        # 设置占位符$PictureIsHere$
        text_doc = aw.Document(docx_path)
        doc_path = os.path.dirname(
            docx_path)+"/"+os.path.basename(docx_path).replace(".docx", ".doc")
        text_doc.save(doc_path, aw.SaveFormat.DOC)

        docx_text = universal_textract(
            doc_path)
        docx_text = docx_text.replace("[pic]", "$PictureIsHere$")
        docx_text = docx_text.replace(
            "Evaluation Only. Created with Aspose.Words. Copyright 2003-2023  Aspose  Pty", "")
        docx_text = docx_text.replace(
            "Ltd.", "")

        flag = True
        logger.info(TAG+"docx_file_info_extract(): 成功设置图片占位符")
        logger.info(TAG+"docx_file_info_extract()-文本信息:")
        logger.info(docx_text)

    except Exception as e:
        logger.info(TAG+"docx_file_info_extract(): 图片占位符设置失败")
        logger.error(e)
        docx_text = " "
        for paragraph in target_docx.paragraphs:
            if "Evaluation Only. Created with Aspose.Words. Copyright 2003-2023  Aspose" not in paragraph.text and "Ltd." not in paragraph.text:
                docx_text += (paragraph.text + "\n")
        logger.info(TAG+"docx_file_info_extract()-文本信息:")
        logger.info(docx_text)

    try:
        # 保存照片
        os.makedirs(image_dir, exist_ok=True)
        dict_rel = target_docx.part.rels
        count = 1
        for rel in dict_rel:
            rel = dict_rel[rel]
            if "image" in rel.target_ref:
                if not os.path.exists(image_dir):
                    os.makedirs(image_dir)
                img_name = re.findall("/(.*)", rel.target_ref)[0]
                word_name = os.path.splitext(os.path.basename(docx_path))[0]
                img_name = image_dir+"/" + \
                    f'{os.path.basename(doc_path)}-image{count}.png'
                with open(img_name, "wb") as f:
                    f.write(rel.target_part.blob)
                count += 1

        # 解析图片信息
        image_all_text_res = []
        if globalVar.flag_list[0] == True:
            image_all_text = ocr_table_batch(image_dir)
            # image_all_text = ocr_batch_textract(image_dir)
            logger.info(TAG+"docx_file_info_extract()-image_all_text: ")
            for row in image_all_text:
                # logger.info(row)
                try:
                    single_result = " ".join(
                        [element for sublist in row[1:] for element in sublist])
                    image_all_text_res.append(single_result)
                except IndexError as e:
                    # 处理 IndexError 异常
                    logger.error(e)
            logger.info(TAG+"docx_file_info_extract()-图片文本信息: ")
            logger.info(image_all_text_res)
        else:
            logger.info(TAG+"docx_file_info_extract()-不处理文件内部的图片!!")

        if flag:
            docx_text_list = re.split(r'(\$PictureIsHere\$)', docx_text)
            res = replace_picture_texts(docx_text_list, image_all_text_res)
            res = "\n".join(res)
            return res
        else:
            res = docx_text+"\n".join(image_all_text_res)
            return res

    except Exception as e:
        logger.error(e)
        logger.error(TAG+"docx_file_info_extract(): 图片信息提取失败, 只返回文本信息")
        return docx_text


# 2-1.解析.ppt/.dps/.pptx
def ppts_file(file_path, type):
    time = datetime.now().strftime("%Y%m%d%H%M%S%f")
    if type == ".ppt":
        logger.info(TAG+"ppts_file(): .ppt")
        ppt_file_dir = "../workspace/office/ppt/"
        result_image_path = "../workspace/image/office/ppt"
        pptx_path = ppt_file_dir + \
            time+"_"+file_path.replace(".ppt", ".pptx").split("/")[-1]
    elif type == ".dps":
        logger.info(TAG+"ppts_file(): .dps")
        ppt_file_dir = "../workspace/wps/dps/"
        result_image_path = "../workspace/image/wps/dps"
        pptx_path = ppt_file_dir + \
            time+"_"+file_path.replace(".dps", ".pptx").split("/")[-1]
    elif type == ".pptx":
        logger.info(TAG+"ppts_file(): .pptx")
        ppt_file_dir = "../workspace/wps/pptx/"
        result_image_path = "../workspace/image/office/pptx"
        ppt_pptx_name = time+"_" + \
            file_path.replace(".dps", ".pptx").split("/")[-1]
        return pptx_file_info_extract(file_path, result_image_path, ppt_pptx_name)
    else:
        return ""

    os.makedirs(ppt_file_dir, exist_ok=True)
    logger.info(TAG+"ppt_file(): "+pptx_path)
    os.makedirs(result_image_path, exist_ok=True)
    ppt_pptx_name = time+"_"+file_path.split("/")[-1]

    try:
        with slides.Presentation(file_path) as presentation:
            presentation.save(pptx_path, slides.export.SaveFormat.PPTX)
        return pptx_file_info_extract(pptx_path, result_image_path, ppt_pptx_name)
    except Exception as e:
        logger.error(e)
        return ""


# 2-2.提取.pptx中的文本和图片文本信息
def pptx_file_info_extract(pptx_path, result_image_path, ppt_pptx_name):

    logger.info(TAG+"ppts_file(): " + pptx_path + " " +
                result_image_path + " "+ppt_pptx_name)

    try:

        presentation = Presentation(pptx_path)
        slide_text = []
        image_count = 1
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:

                # 提取文本内容
                if hasattr(shape, "text"):
                    # 去除水印文字
                    if ("Evaluation only.\nCreated with Aspose.Slides for .NET Standard" not in shape.text):
                        slide_text.append(shape.text)

                # 提取图片
                if shape.shape_type == 13:  # 13 表示图片

                    slide_text.append("$PictureIsHere$")
                    # 提取图片
                    image_dir = f"{result_image_path}/{ppt_pptx_name}/"
                    os.makedirs(image_dir, exist_ok=True)

                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_filename = image_dir + \
                        f"{ppt_pptx_name}_image{image_count}.png"

                    with open(image_filename, "wb") as img_file:
                        img_file.write(image_bytes)

                    image_count += 1

        logger.info(TAG+"pptx_file_info_extract()-文本信息:")
        logger.info(slide_text)

        text_all = " "

        # 解析图片信息
        image_all_text_res = []
        if globalVar.flag_list[0] == True:
            image_folder_path = f"{result_image_path}/{ppt_pptx_name}/"
            image_all_text = ocr_table_batch(image_folder_path)
            # image_all_text = ocr_batch_textract(image_folder_path)
            logger.info(TAG+"ppt_file(): image_all_text: ")
            for row in image_all_text:
                # logger.info(row)
                try:
                    single_result = " ".join(
                        [element for sublist in row[1:] for element in sublist])
                    image_all_text_res.append(single_result)
                except IndexError as e:
                    # 处理 IndexError 异常
                    logger.error(e)
            logger.info(TAG+"pptx_file_info_extract(): 图片文本信息:")
            logger.info(image_all_text_res)
        else:
            logger.info(TAG+"pptx_file_info_extract(): 不处理文件内部的图片!!")

        # 将图片文本信息按顺序插入到文本内部: 文本-文本-图片(如果存在)-文本
        res = replace_picture_texts(slide_text, image_all_text_res)
        res = "\n".join(res)

        return res
    except Exception as e:
        logger.error(e)
        return ""


# 将图片文本信息按顺序插入到文本内部: 文本-文本-图片(如果存在)-文本
def replace_picture_texts(text_list, image_text_list):
    logger.info(TAG+"replace_picture_texts():")

    i = 0  # 用于追踪textlist的索引
    j = 0  # 用于追踪image_all_text_res的索引

    while i < len(text_list):
        if text_list[i] == "$PictureIsHere$":
            if j < len(image_text_list):
                text_list[i] = image_text_list[j]
                i += 1
                j += 1
            else:
                del text_list[i]
        else:
            i += 1

    if j < len(image_text_list):
        text_list.extend(image_text_list[j:])

    return text_list


# 提取.xlsx中的文本
def xlsx_file(file_path):
    xlsx = pd.ExcelFile(file_path)
    sheet_names = xlsx.sheet_names
    # for sheet_name in sheet_names:
    #     data = xlsx.parse(sheet_name)
    xlsx_queue = queue.Queue()
    for name in sheet_names:
        xlsx_queue.put(XlsxDevider(xlsx.parse(name, header=None)))
    # xlsx_queue.put(XlsxDevider(xlsx.parse(sheet_names[5],header=None)))
    res = []
    while not xlsx_queue.empty():
        xlsxDevider = xlsx_queue.get()
        que_add = XlsxDevider.process_xlsx(xlsxDevider)
        if xlsxDevider.check_Pass():
            res.append(xlsxDevider.extract_sensitive_xlsx())
        else:
            while not que_add.empty():
                xlsx_queue.put(que_add.get())
    return res


# 对xlsx中的时间进行格式化
def handle_datetime(obj):
    if isinstance(obj, datetime):
        return obj.strftime('%Y-%m-%d %H:%M:%S')
    return None
