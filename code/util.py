from pptx import Presentation

def convert_ppt_to_pptx(ppt_path, pptx_path):
    # 使用 python-pptx 打开 .ppt 文件
    presentation = Presentation(ppt_path)

    # 创建一个空白的 .pptx 文件
    pptx = Presentation()

    # 遍历 .ppt 中的每个幻灯片
    for slide in presentation.slides:
        # 创建一个新的 .pptx 幻灯片
        pptx_slide = pptx.slides.add_slide(slide.slide_layout)

        # 复制幻灯片中的内容
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_slide.shapes[0].text_frame.text += run.text

    # 保存为 .pptx 文件
    pptx.save(pptx_path)


def read_ppt(ppt_file):
    # 打开PPT文件
    with open(ppt_file, "rb") as f:
        ppt_data = f.read()

    # 读取PPT文件内容
    presentation = Presentation(six.BytesIO(ppt_data))

    # 读取每个幻灯片的内容
    slides = []
    for slide in presentation.slides:
        # 读取幻灯片的标题
        title = slide.shapes.title.text if slide.shapes.title else None

        # 读取幻灯片的正文内容
        content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                content.append(shape.text)

        # 将幻灯片的标题和正文内容添加到列表中
        slides.append({'title': title, 'content': content})

    # 返回所有幻灯片的内容
    return slides

# 调用read_ppt函数读取PPT文件
ppt_file_path = '/home/sakucy/networkCopitation/2023/data/office/学生信息管理系统使用介绍.ppt'
# ppt_content = read_ppt(ppt_file_path)
# 示例用法
pptx_file_path = '/home/sakucy/networkCopitation/2023/workspace/1.pptx'

convert_ppt_to_pptx(ppt_file_path, pptx_file_path)



# # 输出PPT内容
# for slide in ppt_content:
#     print("Slide Title:", slide['title'])
#     print("Slide Content:", slide['content'])
#     print()